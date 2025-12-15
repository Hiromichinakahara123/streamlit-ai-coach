import streamlit as st
import sqlite3
import pandas as pd
from datetime import datetime
from zoneinfo import ZoneInfo
import os
import re
import json
import io
import google.generativeai as genai

# ---------- File parsing ----------
import pypdf
from docx import Document
from pptx import Presentation

# ---------- Gemini ----------
import google.generativeai as genai


# =====================================================
# DB
# =====================================================

DB_FILE = "pk_study_log.db"

def init_db():
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            timestamp TEXT,
            topic TEXT,
            is_correct INTEGER
        )
    """)
    conn.commit()
    conn.close()

def log_result(topic, is_correct):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()

    now_jst = datetime.now(ZoneInfo("Asia/Tokyo"))

    c.execute(
        "INSERT INTO logs (timestamp, topic, is_correct) VALUES (?, ?, ?)",
        (now_jst.strftime("%Y-%m-%d %H:%M:%S"), topic, int(is_correct))
    )

    conn.commit()
    conn.close()


def get_stats():
    conn = sqlite3.connect(DB_FILE)
    df = pd.read_sql("SELECT * FROM logs", conn)
    conn.close()
    return df


# =====================================================
# Gemini
# =====================================================

def configure_gemini():
    api_key = st.secrets.get("GEMINI_API_KEY") or os.getenv("GEMINI_API_KEY")
    if not api_key:
        st.error("❌ GEMINI_API_KEY が設定されていません")
        return False
    genai.configure(api_key=api_key)
    return True


# =====================================================
# File extraction
# =====================================================

def extract_from_pdf(data):
    reader = pypdf.PdfReader(io.BytesIO(data))
    texts = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            texts.append(f"【ページ {i+1}】\n{text}")
    return "\n\n".join(texts)

def extract_from_docx(data):
    doc = Document(io.BytesIO(data))
    texts = []
    for p in doc.paragraphs:
        if p.style.name.startswith("Heading"):
            texts.append(f"\n## {p.text}\n")
        else:
            texts.append(p.text)
    return "\n".join(texts)

def extract_from_xlsx(data):
    xl = pd.ExcelFile(io.BytesIO(data))
    texts = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        texts.append(f"\n## シート: {sheet}\n")
        texts.append(df.to_csv(index=False))
    return "\n".join(texts)

def extract_from_pptx(data):
    prs = Presentation(io.BytesIO(data))
    texts = []
    for i, slide in enumerate(prs.slides):
        texts.append(f"\n## スライド {i+1}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text(uploaded_file):
    data = uploaded_file.read()
    ext = uploaded_file.name.split(".")[-1].lower()

    if ext == "pdf":
        return extract_from_pdf(data)
    if ext == "docx":
        return extract_from_docx(data)
    if ext == "xlsx":
        return extract_from_xlsx(data)
    if ext == "pptx":
        return extract_from_pptx(data)

    raise ValueError("未対応のファイル形式です")


# =====================================================
# AI problem generation
# =====================================================

def safe_json_load(text: str):
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```.*?\n", "", text)
        text = text.rstrip("`").strip()

    # 最初の [ から最後の ] を抽出
    start = text.find("[")
    end = text.rfind("]")
    if start == -1 or end == -1:
        raise ValueError("JSON配列が見つかりません")

    json_text = text[start:end + 1]

    try:
        return json.loads(json_text)
    except json.JSONDecodeError as e:
        raise ValueError(
            f"JSON解析失敗: {e}\n\n--- Gemini出力 ---\n{text}"
        )

def generate_ai_problems(text, n=5):
    model = genai.GenerativeModel("gemini-flash-latest")

    system_prompt = """
あなたは薬剤師国家試験対策問題を作成する教育AIです。

【厳守事項】
・提供資料の内容のみから作問する
・薬剤師国家試験形式（5択単一選択）とする
・正解は必ず1つ
・誤りの選択肢は知識不足で選びやすいものにする
・JSONのみ出力
・JSONのキーや値に改行を含めない
・choicesの各選択肢は1文で完結させる
・説明文は100文字以内
・数式は LaTeX や $ 記法を使わず、すべて文章または通常の記号で書く
・バックスラッシュ（\）を一切使用しない
"""

    prompt = f"""
以下の資料から {n} 問の五肢択一問題を作成してください。

【重要】
・各問題に必ず「topic（分野名）」を付ける
・topicは薬剤師国家試験の科目・領域名で簡潔に書く
  （例：薬物動態学、製剤学、物理薬剤学、薬理学 など）

出力形式:
[
  {{
    "topic": "...",
    "question": "...",
    "choices": {{
      "A": "...",
      "B": "...",
      "C": "...",
      "D": "...",
      "E": "..."
    }},
    "correct": "A",
    "explanation": "..."
  }}
]

資料:
{text[:3000]}
"""

    response = model.generate_content(
        [system_prompt, prompt],
        generation_config={"temperature": 0.2,"response_mime_type": "application/json"}
    )

    return safe_json_load(response.text)


def get_ai_coaching_message(df):
    if df.empty:
        return "まだ学習履歴がありません。"

    # 分野別統計
    stats = df.groupby("topic").agg(
        正解数=("is_correct", "sum"),
        回答数=("id", "count")
    )
    stats["正答率"] = stats["正解数"] / stats["回答数"]
    stats_csv = stats.sort_values("正答率").to_csv()

    model = genai.GenerativeModel("gemini-flash-latest")

    prompt = f"""
あなたは【薬学教育・国家試験指導を専門とする大学教員】です。以下は、ある学生の演習結果（分野別）です。
【分野別成績】
{stats_csv}
この結果から、
① 学問的に理解が不十分と考えられる概念
② 学生が陥りやすい誤解の内容
③ それを克服するための具体的学習方法
をそれぞれ明確に書いてください。

【重要】
・前置きや挨拶文は禁止
・すぐに分析から書き始める
"""

    try:
        response = model.generate_content(
            prompt,
            generation_config={
                "temperature": 0.2,
                "max_output_tokens": 1000
            }
        )
        return response.text

    except Exception as e:
        return f"❌ AIコーチング生成エラー: {e}"


# =====================================================
# UI
# =====================================================

def main():
    st.set_page_config("AIコーチング学習アプリ", layout="centered")
    st.title("📚 AIコーチング学習アプリ")

    init_db()
    if not configure_gemini():
        return

    if "text" not in st.session_state:
        st.session_state.text = None
    if "problems" not in st.session_state:
        st.session_state.problems = []
    if "idx" not in st.session_state:
        st.session_state.idx = 0
    if "answered" not in st.session_state:
        st.session_state.answered = False

    tab1, tab2, tab3 = st.tabs(["資料", "問題演習", "コーチング"])

    # ---------- 資料 ----------
    with tab1:
        file = st.file_uploader(
            "資料をアップロード",
            type=["pdf", "docx", "xlsx", "pptx"]
        )
        if file:
            with st.spinner("資料解析中..."):
                st.session_state.text = extract_text(file)
            st.success("資料を読み込みました")

            if st.button("AI問題を生成"):
                try:
                    st.session_state.problems = generate_ai_problems(
                        st.session_state.text
                    )
                    st.session_state.idx = 0
                    st.success("問題を生成しました")
                    st.rerun()

                except Exception as e:
                    st.error("❌ 問題生成に失敗しました")
                    st.exception(e)

    # ---------- 問題 ----------
    with tab2:
        if not st.session_state.problems:
            st.info("問題がまだありません")
            return

    # --- 全問終了 ---
        if st.session_state.idx >= len(st.session_state.problems):
            st.success("🎉 すべての問題が終了しました！")

            df = get_stats()
            correct = df["is_correct"].sum() if not df.empty else 0
            st.write(f"正解数: {correct} / {len(st.session_state.problems)}")

            if st.button("もう一度最初から"):
                st.session_state.idx = 0
                st.session_state.answered = False
                st.rerun()
            return

        p = st.session_state.problems[st.session_state.idx]
        st.subheader(f"問題 {st.session_state.idx + 1}")
        st.markdown(p["question"])

        choice = st.radio(
            "選択肢",
            options=list(p["choices"].keys()),
            format_func=lambda x: f"{x}: {p['choices'][x]}",
            key=f"choice_{st.session_state.idx}"
        )

        # --- 解答する ---
        if not st.session_state.answered:
            if st.button("解答する"):
                st.session_state.answered = True
                st.session_state.is_correct = (choice == p["correct"])
                topic = p.get("topic", "未分類")
                log_result(topic, st.session_state.is_correct)



        # --- 解答後表示 ---
        if st.session_state.answered:
            if st.session_state.is_correct:
                st.success("正解です 🎉")
            else:
                st.error(f"不正解です。正解は {p['correct']} です。")

            st.markdown("### 解説")
            st.markdown(p["explanation"])

            # --- 次の問題へ ---
            if st.button("次の問題へ"):
                st.session_state.idx += 1
                st.session_state.answered = False
                st.rerun()






    # ---------- コーチング ----------
    with tab3:
        df = get_stats()
        if df.empty:
            st.info("学習履歴がありません")
        else:
            st.subheader("分野別 正答率")
            stats = df.groupby("topic").agg(
                正解数=("is_correct", "sum"),
                回答数=("id", "count")
            )
            stats["正答率"] = stats["正解数"] / stats["回答数"]
            st.dataframe(stats, width="stretch")

            if st.button("AIコーチングを更新"):
                with st.spinner("分析中..."):
                    msg = get_ai_coaching_message(df)
                st.info(msg)


if __name__ == "__main__":
    main()































